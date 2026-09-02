from pathlib import Path
from typing import List
from pptx import Presentation
from src.models.document_models import DocumentModel, ParagraphModel, ImageModel, SourceLocation
from src.utils.logger import get_logger

logger = get_logger("pptx_parser")


def parse_pptx(path: Path) -> DocumentModel:
    prs = Presentation(str(path))
    paragraphs: List[ParagraphModel] = []
    images: List[ImageModel] = []

    for i, slide in enumerate(prs.slides, start=1):
        text_chunks = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        text_chunks.append(txt)
                # images: placeholder note
                if shape.shape_type == 13:  # picture
                    loc = SourceLocation(filename=str(path), slide_number=i)
                    images.append(ImageModel(filename=str(path.name), description="embedded image", location=loc))
            except Exception:
                logger.exception("Error reading shape on slide %s", i)
        if text_chunks:
            loc = SourceLocation(filename=str(path), slide_number=i)
            paragraphs.append(ParagraphModel(text="\n".join(text_chunks), location=loc))

    document = DocumentModel(
        filename=str(path.name),
        document_type="pptx",
        paragraphs=paragraphs,
        images=images,
        metadata={},
        source_locations=[{"filename": str(path)}],
    )
    return document
