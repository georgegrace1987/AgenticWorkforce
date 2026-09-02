import tempfile
from pathlib import Path
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
import pytest

from src.document_processing.document_loader import load_documents
from src.document_processing.document_normalizer import normalize_documents
from src.models.document_models import DocumentModel


def create_sample_docx(path: Path):
    doc = DocxDocument()
    doc.add_paragraph("This is a sample requirement paragraph.")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Header1"
    table.rows[0].cells[1].text = "Header2"
    table.rows[1].cells[0].text = "A"
    table.rows[1].cells[1].text = "B"
    doc.save(str(path))


def create_sample_excel(path: Path):
    df = pd.DataFrame({"col1": [1, 2], "col2": ["x", "y"]})
    df.to_excel(str(path), index=False)


def create_sample_pptx(path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title = shapes.title
    title.text = "Sample Slide"
    prs.save(str(path))


def test_load_and_normalize_documents(tmp_path, monkeypatch):
    files = []

    # TXT
    txt = tmp_path / "sample.txt"
    txt.write_text("Line1\nLine2\n")
    files.append(txt)

    # CSV
    csv = tmp_path / "sample.csv"
    csv.write_text("a,b\n1,2\n")
    files.append(csv)

    # DOCX
    docx = tmp_path / "sample.docx"
    create_sample_docx(docx)
    files.append(docx)

    # Excel
    xlsx = tmp_path / "sample.xlsx"
    create_sample_excel(xlsx)
    files.append(xlsx)

    # PPTX
    pptx = tmp_path / "sample.pptx"
    create_sample_pptx(pptx)
    files.append(pptx)

    # Mock PDF parser to avoid dependency on creating a binary PDF
    from src.document_processing import pdf_parser

    def fake_parse_pdf(path):
        return DocumentModel(filename=path.name, document_type="pdf", paragraphs=[], tables=[], images=[], metadata={}, source_locations=[{"filename": str(path)}])

    monkeypatch.setattr(pdf_parser, "parse_pdf", fake_parse_pdf)
    pdf = tmp_path / "sample.pdf"
    pdf.write_text("fake pdf content")
    files.append(pdf)

    docs = load_documents([str(p) for p in files])
    assert isinstance(docs, list)
    assert len(docs) >= 6

    normalized = normalize_documents(docs)
    assert isinstance(normalized, list)
    # Ensure each normalized doc has document_type set
    for d in normalized:
        assert hasattr(d, "document_type")


if __name__ == "__main__":
    pytest.main([__file__])
