import csv
from io import BytesIO, StringIO
from pathlib import Path

from src.document_processing.docx_loader import extract_docx_text


SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".xlsx", ".csv", ".txt"}


def extract_document_text(filename: str, content: bytes) -> str:
    """Extract readable requirement text from the supported source formats."""
    extension = Path(filename).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {extension or 'unknown'}")

    if extension == ".docx":
        return extract_docx_text(content)
    if extension == ".pptx":
        from pptx import Presentation

        presentation = Presentation(BytesIO(content))
        parts = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_text = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if slide_text:
                parts.append(f"Slide {slide_number}:\n" + "\n".join(slide_text))
        return "\n\n".join(parts)
    if extension == ".xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
        parts = []
        for worksheet in workbook.worksheets:
            rows = [" | ".join("" if value is None else str(value) for value in row) for row in worksheet.iter_rows(values_only=True)]
            rows = [row for row in rows if row.strip(" |")]
            if rows:
                parts.append(f"Sheet {worksheet.title}:\n" + "\n".join(rows))
        return "\n\n".join(parts)
    if extension == ".csv":
        decoded = content.decode("utf-8-sig")
        return "\n".join(" | ".join(row) for row in csv.reader(StringIO(decoded)))

    return content.decode("utf-8-sig")
