from io import BytesIO

from fastapi.testclient import TestClient
from openpyxl import Workbook
from pptx import Presentation
import pytest

from src.api.main import app, sessions, wizard_service
from config.settings import settings


class NoOpProvider:
    def generate_questions(self, state, message):
        return []


@pytest.fixture(autouse=True)
def use_noop_provider():
    original = wizard_service.provider
    original_mcp = settings.MCP_FILESYSTEM_ENABLED
    wizard_service.provider = NoOpProvider()
    settings.MCP_FILESYSTEM_ENABLED = False
    yield
    wizard_service.provider = original
    settings.MCP_FILESYSTEM_ENABLED = original_mcp


def _upload(filename: str, content: bytes) -> None:
    response = TestClient(app).post(
        "/api/documents",
        params={"session_id": f"upload-{filename}"},
        files={"file": (filename, content)},
    )
    assert response.status_code == 200, response.text
    assert response.json()["requirements"]["business_objective"]


def test_documents_endpoint_accepts_text_and_csv() -> None:
    _upload("requirements.txt", b"Customers need an invoice portal.")
    _upload("requirements.csv", b"Role,Access\nAdmin,Full")


def test_documents_endpoint_accepts_xlsx_and_pptx() -> None:
    workbook = Workbook()
    workbook.active.append(["Role", "Access"])
    workbook.active.append(["Admin", "Full"])
    spreadsheet = BytesIO()
    workbook.save(spreadsheet)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Invoice portal requirements"
    slides = BytesIO()
    presentation.save(slides)

    _upload("requirements.xlsx", spreadsheet.getvalue())
    _upload("requirements.pptx", slides.getvalue())


def test_stop_command_persists_for_follow_up_messages() -> None:
    client = TestClient(app)
    session_id = "stop-questions"
    response = client.post("/api/chat", json={"session_id": session_id, "message": "Build an invoice portal."})
    assert response.status_code == 200
    response = client.post("/api/chat", json={"session_id": session_id, "message": "Stop asking questions"})
    assert response.json()["requirements"]["questions_paused"] is True
    response = client.post("/api/chat", json={"session_id": session_id, "message": "Admins can approve invoices."})
    assert response.json()["open_questions"] == []
    sessions.pop(session_id, None)
